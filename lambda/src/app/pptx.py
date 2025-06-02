from pptx import Presentation
from pptx.util import Inches
import tempfile
import os

def create_pptx_from_content(slides_data: list[str]) -> str:
    prs = Presentation()

    for slide_text in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title, content = slide_text.split("\n", 1)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    # 一時ファイル保存
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        pptx_path = tmp.name
        prs.save(pptx_path)
    return pptx_path
